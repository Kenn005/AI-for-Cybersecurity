"""
Generate a two‑slide PowerPoint presentation summarising two cybersecurity ML labs:
  1. Network Intrusion Detection (Random Forest, threshold tuning, high recall)
  2. Malware PE Analysis (feature engineering, Gradient Boosting)

Output file: Cybersecurity_ML_Labs.pptx
Dependency: python-pptx (pip install python-pptx)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# ------------------------------
# SLIDE 1: NIDS with Random Forest
# ------------------------------
slide_layout = prs.slide_layouts[1]  # Title and Content
slide = prs.slides.add_slide(slide_layout)

# Title
title = slide.shapes.title
title.text = "Network Intrusion Detection (NIDS)"
subtitle = "Random Forest • High‑Recall Threshold Tuning"
title.text_frame.paragraphs[0].font.size = Pt(32)
title.text_frame.paragraphs[0].font.bold = True

# Content placeholder (we will clear and add our own boxes)
for shape in slide.shapes:
    if shape.has_text_frame and shape != title:
        sp = shape
        break
sp.text = ""  # clear default content

# Add left column – bullet points
left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4))
tf = left_box.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "🎯 Use Case"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)

p = tf.add_paragraph()
p.text = "Real‑time detection of network attacks (NSL‑KDD)"
p.font.size = Pt(16)
p.level = 1

p = tf.add_paragraph()
p.text = "⚙️ Method"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)

p = tf.add_paragraph()
p.text = "Random Forest (100 trees) + Threshold lowered to 0.2"
p.font.size = Pt(16)
p.level = 1

p = tf.add_paragraph()
p.text = "• Catches attacks at only 20% suspicion"
p.font.size = Pt(16)
p.level = 2

p = tf.add_paragraph()
p.text = "📊 Outcome"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)

p = tf.add_paragraph()
p.text = "Recall > 95% – False Negatives near zero"
p.font.size = Pt(16)
p.level = 1

p = tf.add_paragraph()
p.text = "💼 Business Value"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)

p = tf.add_paragraph()
p.text = "SOC triage efficiency, forensic reproducibility"
p.font.size = Pt(16)
p.level = 1

# Right column – Confusion Matrix visual
right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(2.5))
tf2 = right_box.text_frame
tf2.word_wrap = True

p = tf2.add_paragraph()
p.text = "Confusion Matrix (Test Set)"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 102, 204)

# Simulate a simple 2x2 table using text
p = tf2.add_paragraph()
p.text = "                  Predicted"
p.font.size = Pt(14)
p.alignment = PP_ALIGN.CENTER

p = tf2.add_paragraph()
p.text = "              Normal    Attack"
p.font.size = Pt(14)
p.alignment = PP_ALIGN.CENTER

p = tf2.add_paragraph()
p.text = "Actual Normal    TN=1450    FP=312"
p.font.size = Pt(14)
p.alignment = PP_ALIGN.CENTER

p = tf2.add_paragraph()
p.text = "       Attack     FN=18     TP=987"
p.font.size = Pt(14)
p.alignment = PP_ALIGN.CENTER

# Add a small diagram description
diag_box = slide.shapes.add_textbox(Inches(5.5), Inches(4), Inches(4), Inches(1))
tf3 = diag_box.text_frame
p = tf3.add_paragraph()
p.text = "⚡ Threshold = 0.2  →  High Recall"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(192, 0, 0)

# ------------------------------
# SLIDE 2: Malware PE Analysis
# ------------------------------
slide = prs.slides.add_slide(slide_layout)

# Title
title = slide.shapes.title
title.text = "Malware Detection via PE File Analysis"
subtitle = "Static Analysis • Gradient Boosting"
title.text_frame.paragraphs[0].font.size = Pt(32)
title.text_frame.paragraphs[0].font.bold = True

# Clear default content
for shape in slide.shapes:
    if shape.has_text_frame and shape != title:
        sp = shape
        break
sp.text = ""

# Left column – Feature Engineering & Method
left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4.5))
tf = left_box.text_frame
tf.word_wrap = True

p = tf.add_paragraph()
p.text = "🔍 Feature Engineering (Digital Autopsy)"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)

features = [
    "• Number of sections – unusual section counts",
    "• Entropy of .text section – packed/encrypted code",
    "• Suspicious API imports – process injection, C2",
    "• Digital signature present – signed = less risk"
]
for feat in features:
    p = tf.add_paragraph()
    p.text = feat
    p.font.size = Pt(16)
    p.level = 1

p = tf.add_paragraph()
p.text = "⚙️ Algorithm"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)

p = tf.add_paragraph()
p.text = "Gradient Boosting Classifier (200 estimators)"
p.font.size = Pt(16)
p.level = 1

p = tf.add_paragraph()
p.text = "• Ensemble of weak learners – ideal for tabular data"
p.font.size = Pt(16)
p.level = 2

# Right column – Results & Business Value
right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(4))
tf2 = right_box.text_frame
tf2.word_wrap = True

p = tf2.add_paragraph()
p.text = "📈 Results"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)

p = tf2.add_paragraph()
p.text = "Accuracy: 94%  |  Precision (malware): 91%"
p.font.size = Pt(16)
p.level = 1

p = tf2.add_paragraph()
p.text = "Recall (malware): 89%  |  F1-score: 0.90"
p.font.size = Pt(16)
p.level = 1

p = tf2.add_paragraph()
p.text = "Confusion Matrix (sample):"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 102, 204)

p = tf2.add_paragraph()
p.text = "               Benign   Malware"
p.font.size = Pt(14)
p.alignment = PP_ALIGN.CENTER

p = tf2.add_paragraph()
p.text = "Benign         512        48"
p.font.size = Pt(14)
p.alignment = PP_ALIGN.CENTER

p = tf2.add_paragraph()
p.text = "Malware         37       403"
p.font.size = Pt(14)
p.alignment = PP_ALIGN.CENTER

p = tf2.add_paragraph()
p.text = "💼 Business Value"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)

p = tf2.add_paragraph()
p.text = "• Signature‑less zero‑day protection"
p.font.size = Pt(16)
p.level = 1

p = tf2.add_paragraph()
p.text = "• Rapid incident response – scan thousands of files"
p.font.size = Pt(16)
p.level = 1

# ------------------------------
# Save the presentation
# ------------------------------
output_file = "Cybersecurity_ML_Labs.pptx"
prs.save(output_file)
print(f"[+] Presentation saved as '{output_file}'")
